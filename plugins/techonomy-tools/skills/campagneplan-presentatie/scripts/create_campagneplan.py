#!/usr/bin/env python3
"""
Techonomy Campagneplan Generator v3 — Professional + Google Slides compatible
Doelgroep slides use real Techonomy targeting icons extracted from the template.

Usage: python create_campagneplan.py config.json [output.pptx]

Config fields:
  client, campaign_name, start_date, end_date, total_budget,
  target_audience, channels, campaign_notes,
  awareness_pct, verkeer_pct, conversie_pct

  Per-platform targeting (optional, defaults applied if omitted):
  meta_geo, meta_gender, meta_age, meta_targeting
  google_geo, google_gender, google_age, google_targeting
  tiktok_geo, tiktok_gender, tiktok_age, tiktok_targeting
"""

import json, sys, os
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree

# ── Colors ────────────────────────────────────────────────────────────────────
DARK_NAVY  = RGBColor(0,   0,   49)
DARK_BLUE  = RGBColor(0,   0,  220)
ORANGE_RED = RGBColor(255, 70,   0)
PINK       = RGBColor(212, 20,  90)
LIGHT_GRAY = RGBColor(230, 238, 239)
MID_GRAY   = RGBColor(200, 208, 210)
WHITE      = RGBColor(255, 255, 255)
_HDR_BG    = RGBColor(1,   0,  49)   # #010031 — table header bg
_KANAAL_BG = RGBColor(22,  0, 220)   # #1600DC — kanaal cell bg

TEMPLATE = "/Users/marijnbransen/Downloads/Techonomy PPT met voorbeelden.pptx"

# Icon assets (extracted from template slides 59/60)
_HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(_HERE, '..', 'assets', 'icons')
ICON = {
    'location':  os.path.join(ASSETS, 'icon_location.png'),
    'gender':    os.path.join(ASSETS, 'icon_gender.png'),
    'age':       os.path.join(ASSETS, 'icon_age.png'),
    'interests': os.path.join(ASSETS, 'icon_interests.png'),
}
LOGO = {
    'facebook':   os.path.join(ASSETS, 'logo_facebook.png'),
    'instagram':  os.path.join(ASSETS, 'logo_instagram.png'),
    'google_ads': os.path.join(ASSETS, 'logo_google_ads.png'),
    'tiktok':     os.path.join(ASSETS, 'logo_tiktok.png'),
    'youtube':    os.path.join(ASSETS, 'logo_youtube.png'),
}
ICON_W = Inches(0.82)
ICON_H = Inches(0.82)

# Layout constants
ML   = Inches(0.84)    # left margin (matches template content x)
MT   = Inches(1.55)    # content top
CW   = Inches(11.65)   # content width
SH   = Inches(7.5)     # slide height


# ── Utilities ─────────────────────────────────────────────────────────────────

def get_layout(prs, name):
    for l in prs.slide_layouts:
        if l.name == name: return l
    print(f'⚠  layout "{name}" not found, using fallback')
    return prs.slide_layouts[0]

def delete_all_slides(prs):
    lst = prs.slides._sldIdLst
    for rId in [s.get(qn('r:id')) for s in list(lst)]:
        try: prs.part.drop_rel(rId)
        except KeyError: pass
    while len(lst): lst.remove(lst[0])

def _force_tf(tf, size, bold=None, color=None, name=None):
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            if bold  is not None: r.font.bold  = bold
            if color is not None: r.font.color.rgb = color
            if name  is not None: r.font.name  = name

def set_title(slide, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            ph.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            _force_tf(ph.text_frame, 26, bold=True, color=DARK_NAVY, name='Orbitron')
            return ph

def tb(slide, left, top, width, height):
    """Create a text box with word_wrap and no auto-size."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    shape.text_frame.word_wrap = True
    shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    return shape

def run(para, text, size=11, bold=False, italic=False,
        color=DARK_NAVY, name='Montserrat'):
    r = para.add_run()
    r.text = text
    r.font.name = name; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return r

def add_para(tf, text, size=11, bold=False, italic=False,
             color=DARK_NAVY, name='Montserrat',
             align=PP_ALIGN.LEFT, before=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if before: p.space_before = Pt(before)
    run(p, text, size=size, bold=bold, italic=italic, color=color, name=name)
    return p

def arrow_bullet(tf, text, size=10, space=5):
    p = tf.add_paragraph(); p.space_before = Pt(space)
    run(p, '→  ', size=size, color=PINK)
    run(p, text, size=size, color=DARK_NAVY)

def section_hdr(tf, text, size=12, before=0, first=False):
    add_para(tf, text, size=size, bold=True, color=DARK_NAVY, before=before, first=first)

def accent_line(slide, color=ORANGE_RED, thickness=Pt(3)):
    """Thin colored line under the title area."""
    line = slide.shapes.add_shape(1, ML, Inches(1.35), CW, Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()

def _cell_bg(cell, color):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr')).set('val', str(color).lower())

def _cell_write(cell, text, size=9, bold=False, color=DARK_NAVY, align=PP_ALIGN.LEFT, name='Montserrat'):
    tf = cell.text_frame; p = tf.paragraphs[0]; p.alignment = align
    if p.runs: r = p.runs[0]
    else: r = p.add_run()
    r.text = str(text); r.font.name = name
    r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color

def add_picture_safe(slide, path, left, top, width, height):
    """Add picture only if the file exists."""
    if os.path.exists(path):
        return slide.shapes.add_picture(path, left, top, width, height)
    else:
        print(f'  ⚠  icon not found: {path}')
        return None


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_cover(prs, client, campaign_name, today):
    sl = prs.slides.add_slide(get_layout(prs, 'Titeldia met afbeeldingsveld'))
    for ph in sl.placeholders:
        idx = ph.placeholder_format.idx
        tf = ph.text_frame; tf.auto_size = MSO_AUTO_SIZE.NONE
        if idx == 0:
            ph.text = f'Campagneplan: {campaign_name}'
            _force_tf(tf, 28, bold=True, color=DARK_NAVY, name='Orbitron')
        elif idx == 1:
            ph.text = f'{client} × Techonomy  |  {today}'
            _force_tf(tf, 14, color=DARK_NAVY, name='Montserrat')

def slide_toc(prs, chapters):
    sl = prs.slides.add_slide(get_layout(prs, 'Alleen titel - 1'))
    set_title(sl, 'Inhoudsopgave')
    accent_line(sl, DARK_BLUE)
    box = tb(sl, ML, MT, Inches(8), Inches(4.8))
    tf = box.text_frame
    for i, ch in enumerate(chapters):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0 if i == 0 else 14)
        run(p, f'{str(i+1).zfill(2)}  ', size=20, bold=True,
            color=ORANGE_RED, name='Orbitron')
        run(p, ch, size=20, color=DARK_NAVY, name='Montserrat')

def slide_chapter(prs, number, title):
    sl = prs.slides.add_slide(get_layout(prs, 'Hoofdstuk met afbeeldingsveld'))
    for ph in sl.placeholders:
        idx = ph.placeholder_format.idx
        tf = ph.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE; tf.word_wrap = True
        if idx == 0:
            ph.text = title
            _force_tf(tf, 38, bold=True, color=WHITE, name='Orbitron')
        elif idx == 1:
            ph.text = str(number).zfill(2)
            _force_tf(tf, 52, bold=True, color=WHITE, name='Orbitron')
        else:
            try: ph.text = ''
            except: pass

def slide_briefing(prs, client, audience, notes, period, budget):
    sl = prs.slides.add_slide(get_layout(prs, 'Alleen titel - 1'))
    set_title(sl, 'Briefing & aanleiding')
    accent_line(sl)
    box = tb(sl, ML, MT, CW, Inches(5.5))
    tf = box.text_frame
    rows = [
        ('Aanleiding',       notes or f'Paid advertising campagne voor {client}.'),
        ('Doelstelling',     f'Maximale campagneresultaten realiseren voor {client}.'),
        ('Doelgroep',        audience),
        ('Campagneperiode',  period),
        ('Mediabudget',      f'€ {float(budget):,.0f}'),
    ]
    for i, (hdr, body) in enumerate(rows):
        section_hdr(tf, hdr, size=11, before=0 if i==0 else 12, first=(i==0))
        add_para(tf, body, size=10, before=2)

def slide_timeline(prs, start_date, end_date):
    sl = prs.slides.add_slide(get_layout(prs, 'Alleen titel - 1'))
    set_title(sl, 'Campagne tijdlijn')
    accent_line(sl)

    steps = [
        ('Oplevering\ncampagneplan', ''),
        ('Voorbereiding\n& assets',  ''),
        ('Start paid\ncampagne',     start_date),
        ('Einde\ncampagne',          end_date),
        ('Eind-\nevaluatie',         ''),
    ]
    n = len(steps)

    # Layout geometry
    LINE_Y   = Inches(3.70)   # vertical centre of the horizontal timeline bar
    LINE_H   = Inches(0.06)
    CIRC_D   = Inches(0.40)
    BOX_W    = Inches(1.90)
    BOX_H    = Inches(0.80)
    CONN_H   = Inches(0.38)   # height of the vertical connector stub
    SPACING  = CW / n         # horizontal slot per step

    # 1. Horizontal bar (dark navy)
    bar = sl.shapes.add_shape(1, ML, LINE_Y - LINE_H / 2, CW, LINE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK_NAVY
    bar.line.fill.background()

    for i, (label, dt) in enumerate(steps):
        cx = ML + i * SPACING + SPACING / 2   # horizontal centre of this step

        # Alternate boxes: odd steps above the line, even below
        above = (i % 2 == 0)

        # 2. Vertical connector stub
        if above:
            conn_y = LINE_Y - CONN_H - LINE_H / 2
        else:
            conn_y = LINE_Y + LINE_H / 2
        conn = sl.shapes.add_shape(1, cx - Inches(0.02), conn_y, Inches(0.04), CONN_H)
        conn.fill.solid(); conn.fill.fore_color.rgb = DARK_NAVY
        conn.line.fill.background()

        # 3. Rounded-rect label box (light gray, dark navy text)
        box_x = cx - BOX_W / 2
        if above:
            box_y = conn_y - BOX_H
        else:
            box_y = conn_y + CONN_H
        box = sl.shapes.add_shape(5, box_x, box_y, BOX_W, BOX_H)
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = MID_GRAY
        box.line.width = Pt(0.75)
        tf = box.text_frame; tf.auto_size = MSO_AUTO_SIZE.NONE; tf.word_wrap = True
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run(p, label, size=7, bold=True, color=DARK_NAVY)
        if dt:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            run(p2, dt, size=7, bold=False, color=DARK_BLUE)

        # 4. Circle on the timeline bar (dark blue, white step number)
        circ_x = cx - CIRC_D / 2
        circ_y = LINE_Y - CIRC_D / 2
        circ = sl.shapes.add_shape(9, circ_x, circ_y, CIRC_D, CIRC_D)
        circ.fill.solid(); circ.fill.fore_color.rgb = DARK_BLUE
        circ.line.fill.background()
        ctf = circ.text_frame; ctf.auto_size = MSO_AUTO_SIZE.NONE
        ctf.margin_left = Pt(0); ctf.margin_right = Pt(0)
        ctf.margin_top = Pt(2); ctf.margin_bottom = Pt(0)
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        run(cp, str(i + 1), size=9, bold=True, color=WHITE)

def _channel_logos(channel_name):
    """Return list of LOGO keys for the given channel row string."""
    c = channel_name.lower()
    if 'tiktok' in c:
        return ['tiktok']
    if 'youtube' in c or 'display' in c:
        return ['google_ads', 'youtube']
    if 'search' in c or 'performance max' in c or 'pmax' in c:
        return ['google_ads']
    if 'google' in c:
        return ['google_ads']
    if 'meta' in c or 'facebook' in c or 'instagram' in c:
        return ['facebook', 'instagram']
    return []

def slide_channels(prs, channels, period):
    sl = prs.slides.add_slide(get_layout(prs, 'Alleen titel - 1'))
    set_title(sl, 'Selectie campagnekanalen')
    accent_line(sl)
    rows = []
    for ch in channels:
        cl = ch.lower()
        if 'meta' in cl or 'facebook' in cl:
            rows += [
                ['Awareness', 'Meta – Facebook & Instagram',
                 'Merkbekendheid & bereik vergroten', period],
                ['Verkeer',   'Meta – Traffic campagne',
                 'Websiteverkeer genereren',           period],
                ['Conversie', 'Meta – Conversie campagne',
                 'Conversies & leads realiseren',      period],
            ]
        if 'google' in cl:
            rows += [
                ['Awareness', 'Google Display / YouTube',
                 'Merkbekendheid verbreden',            period],
                ['Verkeer',   'Google Search',
                 'Gerichte websitebezoeken',            period],
                ['Conversie', 'Google Performance Max',
                 'Maximale conversies via AI-bidding',  period],
            ]
        if 'tiktok' in cl:
            rows.append(['Awareness', 'TikTok – In-Feed Ads',
                         'Jongere doelgroep bereiken',  period])
    if not rows:
        rows = [
            ['Awareness', 'Meta / Google Display', 'Bereik & merkbekendheid', period],
            ['Verkeer',   'Meta / Google Search',  'Website traffic',         period],
            ['Conversie', 'Meta / Google PMax',    'Conversies',              period],
        ]

    headers  = ['Fase', 'Kanaal', 'Doelstelling', 'Looptijd']
    col_w    = [Inches(1.4), Inches(3.3), Inches(5.5), Inches(1.45)]
    ROW_H    = Inches(0.52)   # fixed row height — needed for logo positioning
    n_r      = len(rows) + 1
    ts = sl.shapes.add_table(n_r, 4, ML, MT, CW, ROW_H * n_r)
    t = ts.table
    for ci, w in enumerate(col_w): t.columns[ci].width = w
    for ri in range(n_r): t.rows[ri].height = ROW_H

    # Header row — same as mediaplan: _HDR_BG, Orbitron Bold White
    for ci, h in enumerate(headers):
        cell = t.cell(0, ci)
        _cell_bg(cell, _HDR_BG)
        _cell_write(cell, h, size=10, bold=True, color=WHITE, name='Orbitron')

    # Data rows
    fase_color = {'awareness': DARK_BLUE, 'verkeer': ORANGE_RED, 'conversie': PINK}
    LOGO_SZ  = Inches(0.30)
    LOGO_PAD = Inches(0.06)
    kanaal_x = ML + col_w[0]   # x start of Kanaal column

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri + 1, ci)
            if ci == 0:
                fc = fase_color.get(val.lower(), DARK_NAVY)
                _cell_bg(cell, fc)
                _cell_write(cell, val, size=9, bold=True, color=WHITE, name='Orbitron')
            elif ci == 1:
                # Kanaal column — _KANAAL_BG, Orbitron Bold White
                # Keep channel name as text; logos added as overlays below
                _cell_bg(cell, _KANAAL_BG)
                cell.text_frame.paragraphs[0].clear()   # clear default text
                # Pad left margin to leave room for logos
                logos = _channel_logos(val)
                logo_block_w = len(logos) * (LOGO_SZ + LOGO_PAD) + LOGO_PAD
                cell.margin_left  = int(logo_block_w)
                cell.margin_top   = int(Inches(0.06))
                cell.margin_bottom = int(Inches(0.06))
                _cell_write(cell, val, size=8, bold=True, color=WHITE, name='Orbitron')
            else:
                _cell_write(cell, val, size=9, color=DARK_NAVY)

        # Place logos as floating images over the Kanaal cell
        logos = _channel_logos(row[1])
        cell_y = MT + ROW_H * (ri + 1)
        logo_y = cell_y + (ROW_H - LOGO_SZ) / 2
        for li, logo_key in enumerate(logos):
            logo_x = kanaal_x + LOGO_PAD + li * (LOGO_SZ + LOGO_PAD)
            add_picture_safe(sl, LOGO[logo_key], logo_x, logo_y, LOGO_SZ, LOGO_SZ)

def _platform_block(sl, name, subtitle, geo, gender, age, targeting, bx, by, block_w):
    """Draw one platform targeting block: header bar + 4 icons + labels + values."""
    # ── Header bar (2 lines: name bold + subtitle smaller) ───────────────────
    BAR_H = Inches(0.62)
    hdr = sl.shapes.add_textbox(bx, by, block_w, BAR_H)
    hdr.fill.solid(); hdr.fill.fore_color.rgb = DARK_NAVY
    hdr.line.fill.background()
    tf = hdr.text_frame; tf.auto_size = MSO_AUTO_SIZE.NONE; tf.word_wrap = False

    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
    run(p1, f'  {name}', size=14, bold=True, color=WHITE, name='Montserrat')
    if subtitle:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT; p2.space_before = Pt(1)
        run(p2, f'  {subtitle}', size=10, color=LIGHT_GRAY, name='Montserrat')

    # ── 4 icon columns ────────────────────────────────────────────────────────
    spacing  = block_w / 4
    icon_y   = by + BAR_H + Inches(0.22)
    label_y  = icon_y + ICON_H + Inches(0.14)
    value_y  = label_y + Inches(0.32)
    value_h  = Inches(0.75)   # tall enough for 2-line values

    params = [
        ('location',  'Geografie', geo),
        ('gender',    'Geslacht',  gender),
        ('age',       'Leeftijd',  age),
        ('interests', 'Targeting', targeting),
    ]

    for i, (icon_key, label, value) in enumerate(params):
        col_x   = bx + i * spacing
        icon_cx = col_x + (spacing - ICON_W) / 2   # horizontally centered

        add_picture_safe(sl, ICON[icon_key], icon_cx, icon_y, ICON_W, ICON_H)

        lt = sl.shapes.add_textbox(col_x, label_y, spacing, Inches(0.30))
        lt.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        lp = lt.text_frame.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
        run(lp, label, size=9, bold=True, color=DARK_NAVY)

        vt = sl.shapes.add_textbox(col_x, value_y, spacing, value_h)
        vt.text_frame.word_wrap = True; vt.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        vp = vt.text_frame.paragraphs[0]; vp.alignment = PP_ALIGN.CENTER
        run(vp, value, size=10, color=DARK_NAVY)

def slide_audience(prs, audience, channels, cfg):
    sl = prs.slides.add_slide(get_layout(prs, 'Alleen titel - 1'))
    set_title(sl, 'Selectie doelgroepen')
    accent_line(sl)

    has_meta   = any('meta' in c.lower() or 'facebook' in c.lower() for c in channels)
    has_google = any('google' in c.lower() for c in channels)
    has_tiktok = any('tiktok' in c.lower() for c in channels)

    geo_default = 'Nederland'
    gen_default = 'Man / vrouw'
    age_default = '25+'

    # Each platform: (name, subtitle, geo, gender, age, targeting)
    platforms = []
    if has_meta:
        platforms.append(('Meta', 'Facebook & Instagram',
            cfg.get('meta_geo',       geo_default),
            cfg.get('meta_gender',    gen_default),
            cfg.get('meta_age',       age_default),
            cfg.get('meta_targeting', 'Interesses & gedrag,\nlookalike audiences')))
    if has_google:
        platforms.append(('Google Ads', 'Search, Display & PMax',
            cfg.get('google_geo',       geo_default),
            cfg.get('google_gender',    gen_default),
            cfg.get('google_age',       'Alle leeftijden'),
            cfg.get('google_targeting', 'Zoekwoorden,\nin-market audiences')))
    if has_tiktok:
        platforms.append(('TikTok', 'In-Feed Ads',
            cfg.get('tiktok_geo',       geo_default),
            cfg.get('tiktok_gender',    gen_default),
            cfg.get('tiktok_age',       '18–35'),
            cfg.get('tiktok_targeting', 'Interesses,\nbroad audience')))
    if not platforms:
        platforms.append(('Meta / Google', 'Facebook, Instagram & Ads',
            geo_default, gen_default, age_default, audience))

    n = len(platforms)
    gap = Inches(0.5)
    block_w = (CW - gap * (n - 1)) / n

    # Vertically center the block group in the available space
    # block total height = header(0.62) + gap(0.22) + icon(0.82) + gap(0.14) + label(0.32) + value(0.75) = ~2.87"
    BLOCK_H = Inches(2.90)
    available_h = SH - MT - Inches(0.85)   # leave room for footer
    block_y = MT + (available_h - BLOCK_H) / 2

    for i, (name, subtitle, geo, gender, age, targeting) in enumerate(platforms):
        bx = ML + i * (block_w + gap)
        _platform_block(sl, name, subtitle, geo, gender, age, targeting, bx, block_y, block_w)

    # Footer note
    note = tb(sl, ML, Inches(6.45), CW, Inches(0.55))
    arrow_bullet(note.text_frame,
        'Targeting is een startpunt — audiences worden gedurende de campagne '
        'geoptimaliseerd op basis van data.', 9)

def slide_copies(prs, campaign_name, client):
    sl = prs.slides.add_slide(get_layout(prs, 'Alleen titel - 1'))
    set_title(sl, 'Suggestie copies')
    accent_line(sl)

    phases = [
        ('Awareness fase', [
            f'"{campaign_name} – Ontdek wat {client} voor je heeft"',
            f'"{client} nodigt je uit — doe mee en maak het verschil"',
        ]),
        ('Verkeer fase', [
            '"Bekijk alle informatie en reserveer nu jouw plek"',
            '"Meer weten? Klik hier voor het volledige programma"',
        ]),
        ('Conversie fase', [
            '"Bestel nu – beperkt beschikbaar"',
            '"Laatste kans: claim jouw plek vandaag nog"',
        ]),
    ]
    # Two-column layout
    col_w = CW / 2 - Inches(0.3)
    left_phases  = phases[:2]
    right_phases = phases[2:]

    def fill(px, py, phase_list):
        box = tb(sl, px, py, col_w, Inches(5))
        tf = box.text_frame; first = True
        for ph_name, bullets in phase_list:
            section_hdr(tf, ph_name, size=11, before=0 if first else 16, first=first)
            first = False
            for b in bullets: arrow_bullet(tf, b, size=10)

    fill(ML, MT, left_phases)
    fill(ML + col_w + Inches(0.6), MT, right_phases)

    note = tb(sl, ML, Inches(6.45), CW, Inches(0.5))
    add_para(note.text_frame,
        'Tip: test minimaal 3 copy-varianten per fase voor optimale prestaties.',
        size=9, italic=True, first=True)

def slide_assets(prs, channels):
    sl = prs.slides.add_slide(get_layout(prs, 'Alleen titel - 1'))
    set_title(sl, 'Gewenste visuele assets')
    accent_line(sl)

    has_meta   = any('meta' in c.lower() for c in channels)
    has_google = any('google' in c.lower() for c in channels)
    has_tiktok = any('tiktok' in c.lower() for c in channels)

    col_w = CW / 2 - Inches(0.3)
    def fill_col(blocks, cx):
        box = tb(sl, cx, MT, col_w, Inches(5.2))
        tf = box.text_frame; first = True
        for hdr, items in blocks:
            section_hdr(tf, hdr, size=11, before=0 if first else 14, first=first)
            first = False
            for it in items: arrow_bullet(tf, it, size=10)

    left, right = [], []
    if has_meta:
        left.append(('Meta', [
            'Feed afbeelding: 1080 × 1080 px (1:1)',
            'Story / Reels: 1080 × 1920 px (9:16)',
            'Carousel: 1080 × 1080 px, 2–10 kaarten',
            'Video: 15–30 sec, ondertiteld',
        ]))
    if has_google:
        right.append(('Google', [
            'Display banners: 300×250 / 728×90 / 160×600',
            'Responsive Display: afb. + headlines + desc.',
            'YouTube: 15–30 sec skippable video',
        ]))
    if has_tiktok:
        (right if right else left).append(('TikTok', [
            'In-Feed video: 1080 × 1920 px, 9–15 sec',
            'Subtitels en CTA in beeld verwerkt',
        ]))

    fill_col(left  or [('Formaten', ['Lever assets aan per platform-specificaties'])], ML)
    fill_col(right or [],  ML + col_w + Inches(0.6))

    note = tb(sl, ML, Inches(6.45), CW, Inches(0.5))
    add_para(note.text_frame,
        'Aanbeveling: minimaal 3–5 creatieve varianten per fase voor A/B-testing.',
        size=9, italic=True, first=True)

def slide_budget(prs, total, aw_pct, vk_pct, cv_pct):
    sl = prs.slides.add_slide(get_layout(prs, 'Alleen titel - 1'))
    set_title(sl, 'Budget verdeling')
    accent_line(sl)

    aw  = round(total * aw_pct / 100)
    vk  = round(total * vk_pct / 100)
    cv  = round(total * cv_pct / 100)
    fee = round(total * 0.025)

    # Left: budget breakdown
    box = tb(sl, ML, MT, Inches(7.5), Inches(5.3))
    tf = box.text_frame

    section_hdr(tf, 'Verdeling mediabudget', size=13, first=True)
    add_para(tf, '', size=4)  # spacer

    for label, pct, amt in [
        (f'Awareness ({aw_pct}%)', None, aw),
        (f'Verkeer ({vk_pct}%)',    None, vk),
        (f'Conversie ({cv_pct}%)',  None, cv),
    ]:
        p = tf.add_paragraph(); p.space_before = Pt(6)
        run(p, '→  ', size=12, color=PINK)
        run(p, f'{label}:', size=12, color=DARK_NAVY)
        run(p, f'    € {amt:,.0f}', size=12, bold=True, color=DARK_NAVY)

    # Separator
    p_sep = tf.add_paragraph(); p_sep.space_before = Pt(8)
    run(p_sep, '─' * 45, size=7, color=LIGHT_GRAY)

    # Totaal
    p_t = tf.add_paragraph(); p_t.space_before = Pt(4)
    run(p_t, 'Totaal:', size=14, bold=True, color=DARK_NAVY)
    run(p_t, f'    € {total:,.0f}', size=14, bold=True, color=DARK_NAVY)

    # Admin fee + uren
    for lbl, val in [('Admin fee (2,5%):', f'€ {fee:,.0f}'),
                     ('Operationele uren:', 'nader te bepalen')]:
        p = tf.add_paragraph(); p.space_before = Pt(7)
        run(p, f'{lbl}  ', size=11, color=DARK_NAVY)
        run(p, val, size=11, bold=True, color=DARK_NAVY)

    # Grand total
    p_gt = tf.add_paragraph(); p_gt.space_before = Pt(10)
    run(p_gt, 'Totaal campagnebudget:', size=12, bold=True, color=DARK_NAVY)
    run(p_gt, f'  € {total + fee:,.0f}  (excl. btw)', size=12, bold=True, color=DARK_BLUE)

    # Toelichting italic note
    note = tb(sl, ML, Inches(6.3), Inches(10.5), Inches(0.85))
    add_para(note.text_frame,
        'Toelichting: Bovenstaande verdeling is een suggestie. Mocht gedurende de campagne '
        'blijken dat een fase meer / minder budget benodigde heeft, dan kunnen we in overleg '
        'besluiten om te schuiven in budgetten.',
        size=9, italic=True, first=True)

def _mediaplan_table(prs, title, headers, rows, toelichting_lines):
    sl = prs.slides.add_slide(get_layout(prs, 'Alleen titel - 1'))
    set_title(sl, title)
    accent_line(sl, DARK_BLUE)

    n_r = len(rows) + 1
    row_h = Inches(0.46)
    table_h = row_h * n_r
    ts = sl.shapes.add_table(n_r, len(headers), ML, MT, CW, table_h)
    t = ts.table

    # ── Header row: #010031 bg, Orbitron Bold White ───────────────────────────
    for ci, h in enumerate(headers):
        cell = t.cell(0, ci)
        _cell_bg(cell, _HDR_BG)
        tf = cell.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        if p.runs: r = p.runs[0]
        else: r = p.add_run()
        r.text = h; r.font.name = 'Orbitron'; r.font.size = Pt(10)
        r.font.bold = True; r.font.color.rgb = WHITE

    # ── Data rows ─────────────────────────────────────────────────────────────
    for ri, row in enumerate(rows):
        is_total = str(row[0]).lower() == 'totaal'
        for ci, val in enumerate(row):
            cell = t.cell(ri + 1, ci)
            is_kanaal_col = (ci == 0)

            if is_kanaal_col:
                # Kanaal column (incl. Totaal): #1600DC bg + Orbitron white
                _cell_bg(cell, _KANAAL_BG)
                tf = cell.text_frame; p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                if p.runs: r = p.runs[0]
                else: r = p.add_run()
                r.text = str(val); r.font.name = 'Orbitron'
                r.font.size = Pt(9); r.font.bold = True
                r.font.color.rgb = WHITE
            else:
                # Other cells: light bg, Montserrat
                _cell_bg(cell, LIGHT_GRAY if is_total else RGBColor(245, 247, 250))
                _cell_write(cell, val, size=9, bold=is_total, color=DARK_NAVY)

    # ── Opmerking section below table ─────────────────────────────────────────
    note_top = MT + table_h + Inches(0.28)
    note = tb(sl, ML, note_top, CW, SH - note_top - Inches(0.25))
    tf = note.text_frame

    p = tf.paragraphs[0]; p.space_before = Pt(0)
    run(p, 'Opmerking:', size=10, bold=True, color=DARK_NAVY)
    for line in toelichting_lines:
        arrow_bullet(tf, line, size=10)

def slide_awareness(prs, channels, budget):
    rows = []
    active = [c for c in channels if any(x in c.lower() for x in ['meta','google','tiktok'])]
    n = max(len(active), 1)
    for ch in channels:
        cl = ch.lower(); share = round(budget / n)
        if 'meta' in cl:
            rows.append(['Meta – Facebook / Instagram',
                         f'{round(share/5*1000):,}', '€ 5,00', f'€ {share:,}'])
        if 'google' in cl:
            rows.append(['Google Display / YouTube',
                         f'{round(share/3*1000):,}', '€ 3,00', f'€ {share:,}'])
        if 'tiktok' in cl:
            rows.append(['TikTok – In-Feed Ads',
                         f'{round(share/4*1000):,}', '€ 4,00', f'€ {share:,}'])
    if not rows:
        rows = [['Meta – Facebook / Instagram',
                 f'{round(budget*0.5/5*1000):,}', '€ 5,00', f'€ {round(budget*0.5):,}'],
                ['Google Display / YouTube',
                 f'{round(budget*0.5/3*1000):,}', '€ 3,00', f'€ {round(budget*0.5):,}']]
    rows.append(['Totaal', '', '', f'€ {budget:,}'])
    _mediaplan_table(prs, 'Mediaplan – Awareness',
        ['Kanaal', 'Impressies', 'CPM', 'Budget'], rows, [
        'Bereik campagne: Je advertentie wordt getoond aan zoveel mogelijk mensen binnen je doelgroep.',
        'Doel: Merkbekendheid opbouwen en een warm publiek creëren voor latere campagnefasen.',
    ])

def slide_verkeer(prs, channels, budget):
    rows = []
    active = [c for c in channels if any(x in c.lower() for x in ['meta','google'])]
    n = max(len(active), 1)
    for ch in channels:
        cl = ch.lower(); share = round(budget / n)
        if 'meta' in cl:
            cpc, ctr = 0.60, 1.50
            clicks = round(share / cpc)
            rows.append(['Meta – Traffic campagne',
                         f'{round(clicks/(ctr/100)):,}',
                         f'{clicks:,}', f'{ctr:.2f}%', f'€ {cpc:.2f}', f'€ {share:,}'])
        if 'google' in cl:
            cpc, ctr = 1.20, 3.00
            clicks = round(share / cpc)
            rows.append(['Google Search',
                         f'{round(clicks/(ctr/100)):,}',
                         f'{clicks:,}', f'{ctr:.2f}%', f'€ {cpc:.2f}', f'€ {share:,}'])
    if not rows:
        rows = [
            ['Meta – Traffic campagne', '—',
             f'{round(budget*0.5/0.60):,}', '1,50%', '€ 0,60', f'€ {round(budget*0.5):,}'],
            ['Google Search', '—',
             f'{round(budget*0.5/1.20):,}', '3,00%', '€ 1,20', f'€ {round(budget*0.5):,}'],
        ]
    rows.append(['Totaal', '', '', '', '', f'€ {budget:,}'])
    _mediaplan_table(prs, 'Mediaplan – Verkeer',
        ['Kanaal', 'Impressies', 'Klikken', 'CTR', 'CPC', 'Budget'], rows, [
        'Traffic campagne: Gericht op mensen die interesse hebben getoond — retargeting & lookalike.',
        'Doel: Warme bezoekers naar de website brengen voor verdere oriëntatie of directe actie.',
    ])

def slide_conversie(prs, channels, budget):
    rows = []
    active = [c for c in channels if any(x in c.lower() for x in ['meta','google'])]
    n = max(len(active), 1)
    for ch in channels:
        cl = ch.lower(); share = round(budget / n)
        if 'meta' in cl:
            cpa = 15; conv = round(share / cpa)
            rows.append(['Meta – Conversie campagne',
                         f'{round(share/5*1000):,}', f'€ {cpa}', f'€ {share:,}', str(conv)])
        if 'google' in cl:
            cpa = 12; conv = round(share / cpa)
            rows.append(['Google Performance Max',
                         f'{round(share/4*1000):,}', f'€ {cpa}', f'€ {share:,}', str(conv)])
    if not rows:
        rows = [
            ['Meta – Conversie campagne',
             f'{round(budget*0.5/5*1000):,}', '€ 15',
             f'€ {round(budget*0.5):,}', str(round(budget*0.5/15))],
            ['Google Performance Max',
             f'{round(budget*0.5/4*1000):,}', '€ 12',
             f'€ {round(budget*0.5):,}', str(round(budget*0.5/12))],
        ]
    total_cv = sum(int(r[-1]) for r in rows)
    rows.append(['Totaal', '', '', f'€ {budget:,}', str(total_cv)])
    _mediaplan_table(prs, 'Mediaplan – Conversie',
        ['Kanaal', 'Impressies', 'CPA', 'Budget', 'Conversies'], rows, [
        'Conversie campagne: Gericht op mensen klaar om te converteren (aankoop, registratie, etc.).',
        'CPA is richtlijn — werkelijke CPA wordt real-time geoptimaliseerd door het platform.',
    ])


# ── Main ──────────────────────────────────────────────────────────────────────

def create_campagneplan(cfg, output_path):
    prs = Presentation(TEMPLATE)
    print('📂 Template geladen')
    delete_all_slides(prs)

    client   = cfg.get('client', 'Klant')
    name     = cfg.get('campaign_name', 'Campagne')
    start    = cfg.get('start_date', '')
    end      = cfg.get('end_date', '')
    budget   = float(cfg.get('total_budget', 0))
    audience = cfg.get('target_audience', '')
    channels = cfg.get('channels', ['Meta', 'Google Ads'])
    notes    = cfg.get('campaign_notes', '')
    aw_pct   = int(cfg.get('awareness_pct',  40))
    vk_pct   = int(cfg.get('verkeer_pct',    30))
    cv_pct   = int(cfg.get('conversie_pct',  30))
    period   = f'{start} – {end}' if start and end else start or end
    today    = date.today().strftime('%d-%m-%Y')
    chapters = ['Inleiding & debrief', 'Campagne aanpak',
                'Budget verdeling', 'Mediaplan']

    slide_cover(prs, client, name, today);       print('  ✓ Cover')
    slide_toc(prs, chapters);                    print('  ✓ Inhoudsopgave')

    slide_chapter(prs, 1, 'Inleiding & debrief')
    slide_briefing(prs, client, audience, notes, period, budget)
    print('  ✓ H1: Inleiding & debrief')

    slide_chapter(prs, 2, 'Campagne aanpak')
    slide_timeline(prs, start, end)
    slide_channels(prs, channels, period)
    slide_audience(prs, audience, channels, cfg)
    slide_copies(prs, name, client)
    slide_assets(prs, channels)
    print('  ✓ H2: Campagne aanpak')

    slide_chapter(prs, 3, 'Budget verdeling')
    slide_budget(prs, budget, aw_pct, vk_pct, cv_pct)
    print('  ✓ H3: Budget verdeling')

    slide_chapter(prs, 4, 'Mediaplan')
    slide_awareness(prs, channels, round(budget * aw_pct / 100))
    slide_verkeer(prs,   channels, round(budget * vk_pct / 100))
    slide_conversie(prs, channels, round(budget * cv_pct / 100))
    print('  ✓ H4: Mediaplan')

    prs.slides.add_slide(get_layout(prs, 'Logo'))
    prs.save(output_path)
    print(f'\n✅  {output_path}  ({len(prs.slides)} dia\'s)')

if __name__ == '__main__':
    if len(sys.argv) < 2: print(__doc__); sys.exit(1)
    with open(sys.argv[1]) as f: cfg = json.load(f)
    out = sys.argv[2] if len(sys.argv) > 2 \
        else f"campagneplan_{cfg.get('client','klant').lower().replace(' ','_')}.pptx"
    create_campagneplan(cfg, out)
